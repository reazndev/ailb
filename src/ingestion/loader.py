import os
from docx import Document
from pypdf import PdfReader
from pptx import Presentation
from typing import Optional

def load_docx(file_path: str) -> str:
    """Reads text from a .docx file."""
    try:
        doc = Document(file_path)
        full_text = []
        for para in doc.paragraphs:
            full_text.append(para.text)
        return '\n'.join(full_text)
    except Exception as e:
        print(f"Error reading DOCX {file_path}: {e}")
        return ""

def load_pdf(file_path: str) -> str:
    """Reads text from a .pdf file."""
    try:
        reader = PdfReader(file_path)
        full_text = []
        for page in reader.pages:
            text = page.extract_text()
            if text:
                full_text.append(text)
        return '\n'.join(full_text)
    except Exception as e:
        print(f"Error reading PDF {file_path}: {e}")
        return ""

def load_pptx(file_path: str) -> str:
    """Reads text from a .pptx file."""
    try:
        prs = Presentation(file_path)
        full_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    full_text.append(shape.text)
        return '\n'.join(full_text)
    except Exception as e:
        print(f"Error reading PPTX {file_path}: {e}")
        return ""

def load_file_content(file_path: str) -> str:
    """Dispatches to the correct loader based on file extension."""
    _, ext = os.path.splitext(file_path)
    ext = ext.lower()
    
    if ext == ".docx":
        return load_docx(file_path)
    elif ext == ".pdf":
        return load_pdf(file_path)
    elif ext == ".pptx":
        return load_pptx(file_path)
    elif ext == ".txt" or ext == ".md":
        try:
            with open(file_path, 'r', encoding='utf-8') as f:
                return f.read()
        except Exception as e:
             print(f"Error reading text file {file_path}: {e}")
             return ""
    else:
        return f"[Skipped unsupported file type: {ext}]"
